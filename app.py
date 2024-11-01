import streamlit as st
import toml
import asyncio
import aiohttp
import json
from ibm_watsonx_ai.foundation_models import Model
from ibm_watsonx_ai.metanames import GenTextParamsMetaNames as GenParams
import PyPDF2
import docx
import openpyxl
import pptx
import io
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss
import requests
import time
from pymilvus import Collection, DataType, FieldSchema, CollectionSchema, connections, utility, MilvusClient
import os
from pathlib import Path
from datetime import datetime
from tqdm.auto import tqdm
import json
import hashlib
from elasticsearch import Elasticsearch
import warnings
import urllib3
import tempfile
from PIL import Image
#import pytesseract
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# Page setup
st.set_page_config(layout="wide", page_title="watsonx.ai RAG Chat")
# 전역 상수 정의
VECTOR_SEARCH_FAISS = "FAISS"
VECTOR_SEARCH_MILVUS = "Milvus"
VECTOR_SEARCH_WATSON_DISCOVERY = "watsonx Discovery"


# 세션 상태를 초기화하는 함수입니다.
# 이 함수는 애플리케이션의 다양한 설정과 상태 변수들을 초기화합니다.
def init_session_status():
    # 인증 상태 초기화
    if 'authenticated' not in st.session_state:
        st.session_state.authenticated = False
    
    # 채팅 메시지 초기화
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    
    # API 관련 설정 초기화
    if 'url' not in st.session_state:
        st.session_state.url = None  # 기본값으로 Dallas 설정
    if 'apikey' not in st.session_state:
        st.session_state.apikey = ""
    if 'bearer' not in st.session_state:
        st.session_state.bearer = ""
    
    # 모델 및 프로젝트 설정 초기화
    if 'model_id' not in st.session_state:
        st.session_state.model_id = None
    if 'project_id' not in st.session_state:
        st.session_state.project_id = ""
    if 'space_id' not in st.session_state:
        st.session_state.space_id = ""
    
    # 모델 파라미터 초기화
    if 'decoding_method' not in st.session_state:
        st.session_state.decoding_method = "greedy"
    if 'max_new_tokens' not in st.session_state:
        st.session_state.max_new_tokens = 3000
    if 'min_new_tokens' not in st.session_state:
        st.session_state.min_new_tokens = 1
    if 'temperature' not in st.session_state:
        st.session_state.temperature = 0.7
    if 'top_k' not in st.session_state:
        st.session_state.top_k = 50
    if 'top_p' not in st.session_state:
        st.session_state.top_p = 1.0
    
    # 문서 처리 관련 변수 초기화
    if 'document_vectors' not in st.session_state:
        st.session_state.document_vectors = {}
    if 'document_chunks' not in st.session_state:
        st.session_state.document_chunks = {}
    if 'chunk_size' not in st.session_state:
        st.session_state.chunk_size = 200
    if 'overlap_size' not in st.session_state:
        st.session_state.overlap_size = 50
    
    # 토큰 및 프롬프트 관련 변수 초기화
    if 'token_expiry' not in st.session_state:
        st.session_state.token_expiry = 3600
    if 'user_prompt' not in st.session_state:
        st.session_state.user_prompt = None
    if 'last_full_prompt' not in st.session_state:
        st.session_state.last_full_prompt = ""
    
    # UI 관련 변수 초기화
    if 'expanders' not in st.session_state:
        st.session_state.expanders = []
    
    # 모델 선택 및 검색 방법 초기화
    if 'selected_model' not in st.session_state:
        st.session_state.selected_model = None
    if 'vector_search_method' not in st.session_state:
        st.session_state.vector_search_method = VECTOR_SEARCH_FAISS
    
    # Milvus 관련 설정 초기화
    if 'milvus_config' not in st.session_state:
        st.session_state.milvus_config = None
    if 'milvus_test_success' not in st.session_state:
        st.session_state.milvus_test_success = False
    if 'milvus_collection' not in st.session_state:
        st.session_state.milvus_collection = None
    if 'milvus_connected' not in st.session_state:
        st.session_state.milvus_connected = False
    if 'milvus_connection_tested' not in st.session_state:
        st.session_state.milvus_connection_tested = False    
    # 로그 메시지 초기화
    if 'log_messages' not in st.session_state:
        st.session_state.log_messages = ""
    
    # 기타 설정 초기화
    if 'method' not in st.session_state:
        st.session_state.method = None
    if 'lang' not in st.session_state:
        st.session_state.lang = "ko"
    if 'config' not in st.session_state:
        st.session_state.config = None
    if 'labels' not in st.session_state:
        st.session_state.labels = None
    if 'url_options' not in st.session_state:
        st.session_state.url_options = None
    
    # 처리된 파일 목록 초기화
    if 'processed_files' not in st.session_state:
        st.session_state.processed_files = set()
    
    # 문장 임베딩 모델 초기화
    if 'sentence_transformer' not in st.session_state:
        st.session_state.sentence_transformer = None                                                         
    
    # watsonx Discovery 관련 설정 초기화
    if 'watsonx_discovery_config' not in st.session_state:
        st.session_state.watsonx_discovery_config = {
            'server_url': '',
            'api_key': '',
            'pem_content': '',
            'index_name': '',
            'text_embedding_type': None
        }

    if 'selected_model_info' not in st.session_state:
        st.session_state.selected_model_info = None

    if 'wxd_connection_tested' not in st.session_state:
        st.session_state.wxd_connection_tested = False

    if 'milvus_connection_tested' not in st.session_state:
        st.session_state.milvus_connection_tested = False

    if 'model_info' not in st.session_state:
        st.session_state.model_info = None               

    # FAISS 관련 세션 상태 초기화
    if 'faiss_index' not in st.session_state or st.session_state.faiss_index is None:
        # 초기에는 None으로 설정하고, 첫 벡터가 추가될 때 실제 인덱스를 생성합니다.
        st.session_state.faiss_index = None
        st.session_state.faiss_id_to_file_chunk = {}

    if 'document_vectors' not in st.session_state:
        st.session_state.document_vectors = {}
    if 'document_chunks' not in st.session_state:
        st.session_state.document_chunks = {}
    if 'processed_files' not in st.session_state:
        st.session_state.processed_files = set()

    if 'watsonx_discovery_client' not in st.session_state:
        st.session_state.watsonx_discovery_client = None

    if 'stream_endpoint_url' not in st.session_state:
        st.session_state.stream_endpoint_url = None

    if 'login_mode' not in st.session_state:
        st.session_state.login_mode = 'general'

    return



# 지정된 언어와 키에 대한 메시지를 가져오고, 필요한 경우 포맷합니다.
#
# 이 함수는 다국어 지원을 위해 사용되며, 현재 설정된 언어에 따라 적절한 메시지를 반환합니다.
# 메시지에 변수가 포함된 경우, 제공된 키워드 인자를 사용하여 포맷합니다.
#
# :param message_key: 메시지 셔너리의 키
# :param kwargs: 포맷팅에 사용될 키워드 인자들
# :return: 포맷된 (또는 원본) 메시지 문자열
def get_label(message_key, **kwargs):
    try:
        # 현재 설정된 언어로 메시지 가져오기
        message = st.session_state.labels[st.session_state.lang][message_key]
        
        # kwargs가 제공되었다면 포맷을 시도합니다
        if kwargs:
            try:
                return message.format(**kwargs)
            except KeyError as e:
                # 누락된 키가 있다면 오류 메시지를 반환합니다
                return get_label('format_error', message_key=message_key, missing_key=str(e))
        else:
            # kwargs가 없다면 원본 메시지를 그대로 반환합니다
            return message
    except KeyError:
        # 메시지 키나 언어가 존재하지 않는 경우 오류 메시지를 반환합니다
        return get_label('message_not_found', lang=st.session_state.lang, message_key=message_key)

# 로그 메시지를 추가하고 필요에 따라 화면에 표시합니다.
#
# 이 함수는 애플리케이션의 로깅 기능을 담당하며, 메시지를 로그에 추가하고
# 선택적으로 화면에 성공 또는 오류 메시지로 표시합니다.
#
# :param message: 로그에 추가할 메시지
# :param flag: 메시지 표시 방식을 결정하는 플래그 (None, True, False)
#              None: 로그에만 추가
#              True: 성공 메시지로 화면에 표시
#              False: 오류 메시지로 화면에 표시
def add_log(message, flag=None):
    # 현재 시간을 포함한 로그 메시지 생성
    log_entry = datetime.now().strftime('%Y-%m-%d %H:%M:%S') + ")" + message + "\n"
    
    # flag에 따라 메시지 처리
    if flag is None:
        # 로그에만 추가
        st.session_state.log_messages += log_entry
    elif flag:
        # 성공 메시지로 화면에 표시하고 로그에 추가
        st.success(message)
        st.session_state.log_messages += log_entry
    else:
        # 오류 메시지로 화면에 표시하고 로그에 추가
        st.error(message)
        st.session_state.log_messages += log_entry

# 주어진 기본 경로 내에 특정 이름으로 끝나는 디렉토리가 존재하는지 확인합니다.
#
# 이 함수는 주로 모델 캐시 디렉토리의 존재 여부를 확인하는 데 사용됩니다.
#
# :param base_path: 검색을 시작할 기본 디렉토리 경로
# :param target_name: 찾고자 하는 디렉토리 이름의 끝 부분
# :return: 해당 디렉토리가 존재하면 True, 그렇지 않으면 False
def directory_containing_exists(base_path, target_name):
    # 기본 경로를 Path 객체로 변환
    base_path = Path(base_path)
    
    # 기본 경로 내의 모든 항목을 순회
    for item in base_path.iterdir():
        # 항목이 디렉토리이고 이름이 target_name으로 끝나는지 확인
        if item.is_dir() and item.name.endswith(target_name):
            return True
    
    # 조건을 만족하는 디렉토리를 찾지 못한 경우
    return False

# 문장 임베딩 모델들을 로드하고 캐시합니다.
#
# 이 함수는 설정된 임베딩 모델들을 순차적으로 로드하고, 로딩 진행 상황을 
# 사용자에게 시각적으로 보여줍니다. 로드된 모델들은 세션 상태에 저장됩니다.
#
# :return: None

def get_cache_dir():
    cache_dir_template = st.session_state.config['model']['cache_dir']
    cache_dir = os.environ.get('TRANSFORMERS_CACHE')
    if not cache_dir:
        home = os.path.expanduser('~')
        cache_dir = os.path.join(home, '.cache', 'huggingface', 'hub')
    return Path(cache_dir).resolve()

@st.cache_resource
def load_sentence_transformer():
    model_name = st.session_state.selected_model
    
    # 진행 상황을 표시할 컨테이너 생성
    progress_container = st.container()
    subheader = st.subheader(get_label('welcome'))
    with progress_container:
        stitle = st.empty()
        stitle.markdown(get_label("model_loading_progress"), unsafe_allow_html=True)
        progress_bar = st.progress(0)
        status_text = st.empty()
        overall_status = st.empty()

    try:
        # 모델 캐시 디렉토리 설정
        cache_dir = get_cache_dir()
        add_log(f"Cache directory: {cache_dir}")
        print(f"Cache directory: {cache_dir}")
        
        # 캐시 디렉토리가 존재하지 않으면 생성
        cache_dir.mkdir(parents=True, exist_ok=True)
        
        # 모델 로드
        with progress_container:
            if directory_containing_exists(cache_dir, model_name.replace('/', '--')):
                status_text.text(f"{model_name}: {get_label('loading_model', model_name=model_name)}")
            else:
                status_text.text(f"{model_name}: {get_label('downloading_model', model_name=model_name)}")

        # 모델 다운로드 및 로드
        model = SentenceTransformer(model_name, cache_folder=str(cache_dir))
        add_log(f"Model loaded: {model_name}")
        print(f"Model loaded: {model_name}")
        
        # 진행 상황 업데이트
        with progress_container:
            status_text.text(f"{model_name}: {get_label('model_loaded')}")
            progress_bar.progress(1.0)
            overall_status.text(f"{get_label('overall_progress')}: 100%")

        # 로드된 모델을 세션 상태에 저장
        st.session_state.sentence_transformer = model

    except Exception as e:
        add_log(f"Error loading model: {str(e)}", False)
        print(f"Error loading model: {str(e)}")
        status_text.text(f"Error: {str(e)}")
    finally:
        # UI 정리
        progress_container.empty()
        subheader.empty()
        stitle.empty()
        progress_bar.empty()
        overall_status.empty()
        status_text.empty()

    return model


# 애플리케이션의 초기 설정을 수행하는 함수입니다.
#
# 이 함수는 설정 파일을 로드하고, 기본 언어를 설정하며,
# URL 옵션, 모델 ID, 언어 파일 등 애플리케이션의 기본 상태를 초기화합니다.
def initialize():
    # 설정 파일 로드
    with open('config/app.toml', 'r') as f:
        st.session_state.config = toml.load(f)
    
    # 기본 언어 설정
    st.session_state.lang = st.session_state.config['general']['language']
    
    # URL 옵션 및 기본 URL 설정
    st.session_state.url_options = st.session_state.config['url_options']
    st.session_state.url = st.session_state.config['url_options']['Dallas']
    
    # 기본 모델 ID 설정
    st.session_state.model_id = st.session_state.config["model"]["supported_models"][0]
    
    # 기본 임베딩 모델 설정
    st.session_state.selected_model = st.session_state.config["model"]["embedding_models"][0]

    # 언어 파일 로드
    with open(st.session_state.config['general']['language_file'], 'r', encoding='utf-8') as f:
        st.session_state.labels = json.load(f)    
    
    # 사용자 프롬프트 초기화
    st.session_state.user_prompt = get_label("user_prompt_full")
    
    # API 메소드 설정
    st.session_state.method = st.session_state.config["general"]["method"]
    
    # 사용자 정의 CSS 스타일 적용
    st.markdown(st.session_state.config['styles']['css'], unsafe_allow_html=True)

# 제공된 API 키를 사용하여 IBM Cloud IAM 서비스로부터 Bearer 토큰을 얻습니다.
#
# 이 함수는 IBM Cloud IAM 서비스에 API 키를 전송하여 인증을 수행하고,
# 성공 시 Bearer 토큰을 반환합니다. 이 토큰은 IBM 서비스에 대한 후속 API 호출에 사용됩니다.
#
# :param api_key: IBM Cloud API 키
# :return: 성공 시 Bearer 토큰, 실패 시 None
def get_bearer_token(api_key):
    url = "https://iam.cloud.ibm.com/identity/token"
    headers = {
        "Content-Type": "application/x-www-form-urlencoded"
    }
    data = {
        "grant_type": "urn:ibm:params:oauth:grant-type:apikey",
        "apikey": api_key
    }
    
    try:
        # IAM 서비스에 POST 요청 전송
        response = requests.post(url, headers=headers, data=data)
        response.raise_for_status()  # HTTP 오류 발생 시 예외 발생
        
        # 응답에서 토큰 정보 추출
        token_data = response.json()
        st.session_state.bearer = token_data["access_token"]
        st.session_state.token_expiry = time.time() + token_data["expires_in"]

        return st.session_state.bearer
    except requests.RequestException as e:
        # 요청 실패 시 오류 로그 추가
        add_log(get_label('bearer_token_error', status_code=response.status_code), False)
        return None

# Bearer 토큰의 유효성을 확인하고 필요한 경우 갱신합니다.
#
# 이 함수는 현재 토큰의 만료 시간을 확인하고, 만료가 임박한 경우 새로운 토큰을 요청합니다.
# 토큰 갱신 과정에서 발생할 수 있는 오류를 처리하고 적절한 로그를 기록합니다.
def ensure_valid_token():
    # 토큰 만료 5분 전에 갱신을 시도합니다
    if time.time() > st.session_state.token_expiry - 300:
        if st.session_state.apikey:
            with st.spinner(get_label('renewing_token')):
                # 새로운 토큰 요청
                new_token = get_bearer_token(st.session_state.apikey)
                if new_token:
                    add_log(get_label('token_renewal_success'), True)
                else:
                    add_log(get_label('token_renewal_failure'), False)
        else:
            # API 키가 설정되지 않은 경우 경고 메시지 표시
            st.warning(get_label('api_key_not_set'))

# 다양한 형식의 파일에서 텍스트를 추출합니다.
#
# 이 함수는 PDF, DOCX, XLSX/XLS, PPTX 형식의 파일을 지원하며,
# 각 파일 형식에 맞는 라이브러리를 사용하여 텍스트를 추출합니다.
#
# :param file: 파일 객체
# :param file_name: 파일의 이름 (확장자 포함)
# :return: 추출된 텍스트 문자열
def extract_text(file, file_name):
    # 파일 확장자 추출
    file_extension = file_name.split('.')[-1].lower()
    
    try:
        if file_extension == 'pdf':
            return extract_text_from_pdf(file)
        elif file_extension == 'docx':
            return extract_text_from_docx(file)
        elif file_extension in ['xlsx', 'xls']:
            return extract_text_from_excel(file)
        elif file_extension == 'pptx':
            return extract_text_from_powerpoint(file)
        elif file_extension in ['png', 'jpg', 'jpeg', 'gif', 'bmp']:
            return extract_text_from_image(file)
        else:
            add_log(get_label('unsupported_file_format', format=file_extension))
            return ""
    except Exception as e:
        add_log(get_label('text_extraction_error', file_name=file_name, error=str(e)))
        return ""

def extract_text_from_pdf(file):
    try:
        # file이 이미 BytesIO 객체인 경우, 바이트로 변환
        if isinstance(file, io.BytesIO):
            file = file.getvalue()
        
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file))
        text = []
        for page_num, page in enumerate(pdf_reader.pages, 1):
            try:
                page_text = page.extract_text()
                if page_text.strip():
                    text.append(page_text)
                else:
                    # 텍스트 추출 실패 시 이미지에서 텍스트 추출 시도
                    image = extract_image_from_pdf_page(page)
                    if image:
                        text.append(extract_text_from_image(image))
            except Exception as e:
                add_log(get_label('pdf_page_extraction_error', page_num=page_num, error=str(e)))
        return ' '.join(text)
    except Exception as e:
        add_log(get_label('pdf_reading_error', error=str(e)))
        return ""

def extract_image_from_pdf_page(page):
    for image_file_object in page.images:
        return image_file_object.data
    return None

def extract_text_from_docx(file):
    try:
        doc = docx.Document(io.BytesIO(file))
        return ' '.join(paragraph.text for paragraph in doc.paragraphs)
    except Exception as e:
        add_log(get_label('docx_reading_error', error=str(e)))
        return ""

def extract_text_from_excel(file):
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(file), data_only=True)
        text = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text.append(str(cell.value))
        return ' '.join(text)
    except Exception as e:
        add_log(get_label('excel_reading_error', error=str(e)))
        return ""

def extract_text_from_powerpoint(file):
    try:
        prs = pptx.Presentation(io.BytesIO(file))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return ' '.join(text)
    except Exception as e:
        add_log(get_label('powerpoint_reading_error', error=str(e)))
        return ""

def extract_text_from_image(file):
    try:
#        image = Image.open(io.BytesIO(file))
#        text = pytesseract.image_to_string(image) # 여기서는 이미지 처리까지는 하지 않음
        return ""
    except Exception as e:
        add_log(get_label('image_reading_error', error=str(e)))
        return ""
        
# 문서를 지정된 크기의 청크로 분할합니다.
#
# 이 함수는 긴 문서를 더 작은 청크로 나누어 처리하기 쉽게 만듭니다.
# 청크 간에 중복을 허용하여 문맥의 연속성을 유지합니다.
#
# :param document: 분할할 문서 텍스트
# :param chunk_size: 각 청크의 단어 수
# :param overlap_size: 인접한 청크 간 중복되는 단어 수
# :return: 분할된 청크들의 리스트
def split_document_to_chunks(document, chunk_size, overlap_size):
    words = document.split()
    add_log(f"Total words in document: {len(words)}")
    chunks = []
    
    for i in range(0, len(words), chunk_size - overlap_size):
        chunk = ' '.join(words[i:i + chunk_size])
        chunks.append(chunk)
    
    add_log(f"Number of chunks created: {len(chunks)}")
    return chunks


# Milvus 컬렉션을 초기화하고, 필요한 경우 새로 생성합니다.
# 컬렉션에 인덱스를 생성하고 로드합니다.
#
# :param collection_name: 초기화할 컬렉션의 이름
# :param vector_dim: 벡터의 차원
# :return: 초기화된 Milvus 컬렉션 객체
def initialize_milvus_collection(collection_name, vector_dim):
    if not utility.has_collection(collection_name):
        # 새로운 컬렉션 생성
        fields = [
            FieldSchema(name="id", dtype=DataType.INT64, is_primary=True, auto_id=True),
            FieldSchema(name="file_name", dtype=DataType.VARCHAR, max_length=255),
            FieldSchema(name="content", dtype=DataType.VARCHAR, max_length=65535),
            FieldSchema(name="vector", dtype=DataType.FLOAT_VECTOR, dim=vector_dim)
        ]
        schema = CollectionSchema(fields, get_label('collection_description'))
        collection = Collection(collection_name, schema)
        
        # 인덱스 생성
        index_params = {
            "index_type": "IVF_FLAT",
            "metric_type": "L2",
            "params": {"nlist": 1024}
        }
        collection.create_index("vector", index_params)
        add_log(get_label('collection_created', collection_name=collection_name))
    else:
        # 기존 컬렉션 로드
        collection = Collection(collection_name)
        if not collection.has_index():
            # 인덱스가 없는 경우 생성
            index_params = {
                "index_type": "IVF_FLAT",
                "metric_type": "L2",
                "params": {"nlist": 1024}
            }
            collection.create_index("vector", index_params)
            add_log(get_label('index_added', collection_name=collection_name))
    
    # 인덱스 로드
    collection.load()
    add_log(get_label('index_loaded', collection_name=collection_name))
    return collection

# Milvus 데이터베이스에 연결을 시도합니다.
#
# :param milvus_config: Milvus 연결 설정 정보를 담은 딕셔너리
# :return: 연결 성공 시 True, 실패 시 False
def connect_to_milvus(milvus_config):
    try:
        # Milvus 서버에 연결
        connections.connect(
            "default", 
            host=milvus_config["milvus_host"], 
            port=milvus_config["milvus_port"], 
            secure=True, 
            server_pem_path=milvus_config["server_pem_path"], 
            server_name=milvus_config["server_name"], 
            user=milvus_config["user"], 
            password=milvus_config["password"]
        )
        add_log(get_label('milvus_connected'))
        return True
    except Exception as e:
        # 연결 실패 시 오류 로그 추가
        add_log(get_label('milvus_connection_failed_log_with_error', error=str(e)), False)
        return False



#업로드된 파일을 처리하여 텍스트를 추출하고, 청크로 분할한 후 벡터화합니다.
#이 함수는 Streamlit의 캐시 기능을 사용하여 성능을 최적화합니다.

# :param file_content: 파일의 바이너리 내용
# :param file_name: 파일의 이름
# :return: 청크 리스트와 해당 청크들의 벡터 표현

@st.cache_data
def process_uploaded_file_cached(file_content, file_name):
    add_log(f"Processing file: {file_name}")
    text = extract_text(io.BytesIO(file_content), file_name)
    add_log(f"Extracted text length: {len(text)}")
    
    chunks = split_document_to_chunks(text, st.session_state.chunk_size, st.session_state.overlap_size)
    add_log(f"Number of chunks after splitting: {len(chunks)}")
    
    if not chunks:
        add_log(get_label('no_text_extracted', file_name=file_name), False)
        return [], np.array([])
    
    if st.session_state.vector_search_method != VECTOR_SEARCH_WATSON_DISCOVERY:
        chunk_vectors = st.session_state.sentence_transformer.encode(chunks)
        add_log(f"Number of chunk vectors: {len(chunk_vectors)}")
    else:
        chunk_vectors = None
    
    return chunks, chunk_vectors

# 파일의 청크와 벡터를 Milvus 데이터베이스에 삽입합니다.
#
# :param file_name: 삽입할 파일의 이름
# :param chunks: 파일의 텍스트 청크 리스트
# :param vectors: 청크에 대응하는 벡터 리스트
# :param milvus_config: Milvus 연결 및 컬렉션 정보
def ingest_to_milvus(file_name, chunks, vectors):
    milvus_config = st.session_state.milvus_config
    try:
        collection_name = milvus_config["collection_name"]
        
        # 연결 확인
        if not connections.has_connection("default"):
            if not connect_to_milvus(milvus_config):
                raise Exception(get_label('milvus_connection_failed'))

        # 컬렉션 존재 여부 확인 및 생성
        collection = create_collection(vectors, collection_name)
        
        # 컬렉션 로드
        collection.load()

        # 삽입할 엔티티 준비
        entities = [
            [file_name] * len(chunks),  # file_name
            chunks,  # content
            vectors.tolist()  # vector
        ]
        
        # 엔티티 삽입
        insert_result = collection.insert(entities)
        
        # 결과 확인
        if insert_result.insert_count == len(chunks):
            collection.flush()
            add_log(get_label('milvus_ingest_success', file_name=file_name, chunk_count=len(chunks)), True)
        else:
            raise Exception(get_label('milvus_insert_count_mismatch', expected=len(chunks), actual=insert_result.insert_count))

    except Exception as e:
        add_log(get_label('milvus_ingest_error', error=str(e)), False)
        raise  # 상위 레벨에서 처리할 수 있도록 예외를 다시 발생시킵니다.

    finally:
        if 'collection' in locals():
            collection.release()

# Milvus 컬렉션을 생성하고 인덱스를 생성합니다.
#
# :param vectors: 벡터 데이터
# :param collection_name: 컬렉션 이름
# :return: 생성된 Milvus 컬렉션 객체    
def create_collection(vectors, collection_name):
    if not utility.has_collection(collection_name):
        fields = [
                FieldSchema(name="id", dtype=DataType.INT64, is_primary=True, auto_id=True),
                FieldSchema(name="file_name", dtype=DataType.VARCHAR, max_length=255),
                FieldSchema(name="content", dtype=DataType.VARCHAR, max_length=65535),
                FieldSchema(name="vector", dtype=DataType.FLOAT_VECTOR, dim=vectors.shape[1])
            ]
        schema = CollectionSchema(fields, get_label('collection_description'))
        collection = Collection(collection_name, schema)
            
            # 인덱스 생성
        index_params = {
                "index_type": "IVF_FLAT",
                "metric_type": "L2",
                "params": {"nlist": 1024}
            }
        collection.create_index("vector", index_params)
        st.session_state.milvus_collection = collection
    else:
        st.session_state.milvus_collection = Collection(collection_name)
    return st.session_state.milvus_collection

# 로그인 페이지를 구성하고 사용자 인증을 처리합니다.
def login_page():
    # 일반/프롬프트템플릿 선택 (라디오 버튼)
    login_mode = st.radio(
        get_label('select_login_mode'),
        [get_label('login_mode_general'), get_label('login_mode_prompt_template')],
        index=0  # "일반"을 기본값으로 설정
    )
    
    if login_mode == get_label('login_mode_general'):
        login_watsonx_ai_page()
        st.session_state.login_mode = "general"
    else:
        login_watsonx_ai_prompt_template_page()
        st.session_state.login_mode = "prompt_template"
    # 벡터 검색 방법 선택 (라디오 버튼)
    vector_search_method = st.radio(
        get_label('select_vector_search_method'),
        [VECTOR_SEARCH_FAISS, VECTOR_SEARCH_MILVUS, VECTOR_SEARCH_WATSON_DISCOVERY],
        index=0  # FAISS를 기본값으로 설정
    )
    
    st.session_state.vector_search_method = vector_search_method
    
    vector_search_settings()

    test_result = False
    if vector_search_method == VECTOR_SEARCH_MILVUS:
        # Milvus 설정 페이지 호출
        milvus_login()
        test_result = st.session_state.milvus_connection_tested        
    elif vector_search_method == VECTOR_SEARCH_WATSON_DISCOVERY:
        # watsonx Discovery 설정 페이지 호출
        watsonx_discovery_login()
        test_result = st.session_state.wxd_connection_tested
    else:
        st.session_state.milvus_config = None
        st.session_state.watsonx_discovery_config = None
    
    # 로그인 버튼
    if st.button(get_label('enter')):
        # 설정 완료 여부 확인
        check_settings_and_go(vector_search_method, test_result)

def login_watsonx_ai_prompt_template_page():
    st.title(get_label('welcome'))
    st.subheader(get_label('prompt_template_login_title'))

    # Stream endpoint URL 입력
    st.session_state.stream_endpoint_url = st.text_input(
        get_label('stream_endpoint_url'),
        value=st.session_state.get('stream_endpoint_url', ''),
        key='stream_endpoint_url_input'
    )

    # API 키 입력
    api_key = st.text_input(get_label('api_key'), value=st.session_state.apikey, type="password", key="api_key_input")
    if api_key != st.session_state.apikey:
        st.session_state.apikey = api_key
        st.session_state.bearer = ""

    # Bearer 토큰 생성 버튼
    if api_key:
        if st.button(get_label('generate_bearer_token')):
            gen_bearer_token(api_key)

    # Bearer 토큰 표시 (수정 불가)
    st.text_input(get_label('bearer_token'), value=st.session_state.bearer, type="password", disabled=True)

    # 벡터 검색 결과 파라미터와 사용자 입력 파라미터를 나란히 배치
    col1, col2 = st.columns(2)

    with col1:
        st.session_state.vector_search_params = st.text_input(
            get_label('vector_search_params'),
            value=st.session_state.get('vector_search_params', ''),
            key='vector_search_params_input'
        )

    with col2:
        st.session_state.user_input_params = st.text_input(
            get_label('user_input_params'),
            value=st.session_state.get('user_input_params', ''),
            key='user_input_params_input'
        )

# 설정이 완료되었는지 확인하고 인증 성공 여부를 결정합니다.
#
# :param vector_search_method: 벡터 검색 방법
# :param test_result: 연결 테스트 결과
# :return: 설정이 완료되었으면 True, 아니면 False
def check_settings_and_go(vector_search_method, test_result):
    missing_settings = are_settings_complete(vector_search_method, test_result)

    if not missing_settings:
        add_log(get_label('all_settings_complete'), True)
        st.session_state.authenticated = True
        add_log(get_label('auth_success'), True)
        add_log(json.dumps(st.session_state.watsonx_discovery_config, indent=2))
        
        if vector_search_method != VECTOR_SEARCH_WATSON_DISCOVERY:
            with st.spinner(get_label('loading_models_spinner')):
                load_sentence_transformer()          
        
        time.sleep(1.0)
        st.rerun()
    else:
        if vector_search_method != VECTOR_SEARCH_FAISS and not test_result:
            add_log(get_label('connection_failed_log'), False)
        add_log(get_label('missing_settings', settings=", ".join(missing_settings)), False)
        add_log(get_label('fill_all_fields'), False)
    

def watsonx_discovery_login():
    st.session_state.watsonx_discovery_config = watsonx_discovery_config_page()
    if st.button(get_label('test_watsonx_discovery_connection')):
        if st.session_state.wxd_connection_tested and st.session_state.model_info:
            st.success(get_label('watsonx_discovery_connection_success'))
        else:
            st.session_state.wxd_connection_tested, st.session_state.model_info = test_watsonx_discovery_connection(st.session_state.watsonx_discovery_config)
            if st.session_state.wxd_connection_tested and st.session_state.model_info:
                st.success(get_label('watsonx_discovery_connection_success'))
    if st.session_state.wxd_connection_tested and st.session_state.model_info:
        handle_model_selection(st.session_state.model_info)


def milvus_login():
    st.session_state.milvus_config = milvus_config_page()
    if st.button(get_label('test_milvus_connection')):
        st.session_state.milvus_connection_tested = test_milvus_connection(st.session_state.milvus_config)
        if not st.session_state.milvus_connection_tested:
            st.error(get_label('milvus_connection_failed'))



def login_watsonx_ai_page():
    st.title(get_label('welcome'))
    st.subheader(get_label('login_title'))
    st.markdown(get_label('login_instruction'))

    # URL 및 메소드 설정을 위한 컬럼 생성
    col1, col2, col3 = st.columns([1, 3, 3])
    with col1:
        # 위치 선택 드롭다운
        selected_location = st.selectbox(
            get_label('select_location'),
            options=list(st.session_state.url_options.keys()),
            index=list(st.session_state.url_options.values()).index(st.session_state.url)
        )        
        st.session_state.url = st.session_state.url_options[selected_location]
    with col2:
        # URL 표시 (수정 불가)
        st.session_state.url = st.text_input(get_label('url'), value=st.session_state.url, disabled=True)
    with col3:
        # URI 메소드 표시 (수정 불가)
        st.session_state.method = st.text_input(get_label('uri_method'), value=st.session_state.method, disabled=True)
    
    # API 키 입력
    api_key = st.text_input(get_label('api_key'), value=st.session_state.apikey, type="password", key="api_key_input")
    if api_key != st.session_state.apikey:
        st.session_state.apikey = api_key
        st.session_state.bearer = ""

    # Bearer 토큰 생성 버튼
    if api_key:
        if st.button(get_label('generate_bearer_token')):
            gen_bearer_token(api_key)

    # Bearer 토큰 표시 (수정 불가)
    st.text_input(get_label('bearer_token'), value=st.session_state.bearer, type="password", disabled=True)
    
    # 프로젝트 ID 및 스페이스 ID 입력
    st.session_state.project_id = st.text_input(get_label('project_id'), value=st.session_state.project_id)
    st.session_state.space_id = st.text_input(get_label('space_id'), value=st.session_state.space_id, placeholder=get_label('space_id_placeholder'), disabled=True)

def gen_bearer_token(api_key):
    with st.spinner(get_label('getting_bearer_token')):
        bearer_token = get_bearer_token(api_key)
        if bearer_token:
            st.session_state.bearer = bearer_token
            add_log(get_label('bearer_token_success'), True)
        else:
            st.session_state.bearer = ""
            add_log(get_label('bearer_token_failure'), False)

# watsonx Discovery 설정을 위한 페이지를 생성합니다.
# 사용자로부터 watsonx Discovery 연결에 필요한 정보를 입력받습니다.
#
# :return: watsonx Discovery 설정 정보를 담은 딕셔너리

def initialize_watsonx_discovery_config():
    if 'watsonx_discovery_config' not in st.session_state or not st.session_state.watsonx_discovery_config:
        st.session_state.watsonx_discovery_config = {
            'server_url': '',
            'api_key': '',
            'pem_content': '',
            'index_name': '',
            'text_embedding_type': None
        }

def input_watsonx_discovery_details():
    st.session_state.watsonx_discovery_config['server_url'] = st.text_input(
        get_label("watsonx_discovery_server_url"), 
        value=st.session_state.watsonx_discovery_config['server_url']
    )
    
    st.session_state.watsonx_discovery_config['api_key'] = st.text_input(
        get_label("watsonx_discovery_api_key"), 
        type="password", 
        value=st.session_state.watsonx_discovery_config['api_key']
    )
    
    st.session_state.watsonx_discovery_config['index_name'] = st.text_input(
        get_label("watsonx_discovery_index_name"), 
        value=st.session_state.watsonx_discovery_config['index_name']
    )

def handle_pem_file_upload():
    pem_file = st.file_uploader(
        get_label("watsonx_discovery_pem_file"), 
        type=['pem'], 
        accept_multiple_files=False
    )
    
    if pem_file:
        if st.session_state.watsonx_discovery_config.get('pem_content') and os.path.exists(st.session_state.watsonx_discovery_config['pem_content']):
            try:
                os.remove(st.session_state.watsonx_discovery_config['pem_content'])
            except Exception as e:
                add_log(get_label('pem_file_deletion_error', error=str(e)))
        
        with tempfile.NamedTemporaryFile(delete=False, mode='wb', suffix='.pem') as tmp_file:
            tmp_file.write(pem_file.read())
            tmp_pem_path = tmp_file.name
        
        st.session_state.watsonx_discovery_config['pem_content'] = tmp_pem_path
        add_log(get_label('pem_file_uploaded_and_saved'))
    else:
        st.session_state.watsonx_discovery_config['pem_content'] = st.session_state.watsonx_discovery_config.get('pem_content', '')

def watsonx_discovery_config_page():
    initialize_watsonx_discovery_config()
    st.info(get_label('watsonx_discovery_config_info'))
    
    input_watsonx_discovery_details()
    handle_pem_file_upload()

    return st.session_state.watsonx_discovery_config

def extract_embedding_field(data, current_path=[]):
    if isinstance(data, dict):
        for key, value in data.items():
            new_path = current_path + [key]
            if key == 'predicted_value':
                return '.'.join(['ml'] + new_path)
            result = extract_embedding_field(value, new_path)
            if result:
                return result
    return None

def analyze_ml_field(ml_field):
    results = []
    model_info = []

    if 'inference' in ml_field:
        for pipeline_name, pipeline_data in ml_field['inference'].items():
            results.append(get_label('pipeline_ml_field_log', pipeline_name=pipeline_name))
            
            model_id = pipeline_data.get('model_id')
            if model_id:
                embedding_field = extract_embedding_field(pipeline_data, ['inference', pipeline_name])
                
                if embedding_field:
                    predicted_value = pipeline_data.get('predicted_value')
                    
                    if isinstance(predicted_value, dict):
                        embedding_type = "text_expanding"
                        results.append(get_label('elser_model_ml_field_log'))
                        results.append(get_label('model_id_ml_field_log', model_id=model_id))
                        results.append(get_label('keyword_count_ml_field_log', count=len(predicted_value)))
                        results.append(get_label('top_10_keywords_ml_field_log'))
                        sorted_keywords = sorted(predicted_value.items(), key=lambda x: x[1], reverse=True)[:10]
                        for keyword, score in sorted_keywords:
                            results.append(get_label('keyword_score_ml_field_log', keyword=keyword, score=score))
                    
                    elif isinstance(predicted_value, list):
                        embedding_type = "text_embedding"
                        results.append(get_label('embedding_model_ml_field_log'))
                        results.append(get_label('model_id_ml_field_log', model_id=model_id))
                        results.append(get_label('vector_dimension_ml_field_log', dimension=len(predicted_value)))
                        results.append(get_label('vector_sample_ml_field_log'))
                        for value in predicted_value[:5]:
                            results.append(get_label('vector_value_ml_field_log', value=value))
                    
                    else:
                        embedding_type = "unknown"
                        results.append(get_label('unknown_model_type_ml_field_log'))
                        results.append(get_label('model_id_ml_field_log', model_id=model_id))
                        results.append(get_label('predicted_value_type_ml_field_log', value_type=type(predicted_value)))
                    
                    model_info.append({
                        "embedding_type": embedding_type,
                        "model_id": model_id,
                        "embedding_field": embedding_field
                    })

    if not model_info:
        results.append(get_label('no_model_info_ml_field_log'))
        results.append(get_label('ml_field_content_ml_field_log', content=ml_field))

    add_log(get_label('debug_model_info_ml_field_log', model_info=model_info))  # 디버그 출력 추가
    return "\n".join(results), model_info


def create_es_client(config):
    try:
        if config.get('pem_content') and os.path.exists(config['pem_content']):
            # pem_content가 존재하고 파일 경로가 유효한 경우 SSL 인증서 검증 활성화 및 PEM 파일 경로 사용
            es_client = Elasticsearch(
                [config['server_url']],
                api_key=config['api_key'],
                verify_certs=True,
                ca_certs=config['pem_content']
            )
            add_log(get_label('elasticsearch_client_created_with_pem'))
        else:
            # pem_content가 없거나 파일 경로가 유효하지 않은 경우 SSL 인증서 검증 비활성화
            # InsecureRequestWarning 경고 무시
            warnings.filterwarnings("ignore", category=urllib3.exceptions.InsecureRequestWarning)
            
            es_client = Elasticsearch(
                [config['server_url']],
                api_key=config['api_key'],
                verify_certs=False,
            )
            add_log(get_label('elasticsearch_client_created_without_pem'))
        return es_client
    except Exception as e:
        add_log(get_label('elasticsearch_client_creation_error', error=str(e)))
        return None

def ingest_document(es_client, index, file_name, content):
    try:
        doc = {
            "file_name": file_name,
            "content": content,
            "_extract_binary_content": True,
            "_reduce_whitespace": True,
            "_run_ml_inference": True
        }
        result = es_client.index(index=index, body=doc, pipeline=index)
        if result['result'] == 'created':
            add_log(get_label('document_indexing_success', doc_id=result['_id']))
            return result['_id']
        else:
            add_log(get_label('document_indexing_error', error=str(result)))
            return None
    except Exception as e:
        add_log(get_label('document_indexing_exception', error=str(e)))
        return None

def create_and_ingest_test_document(es_client, index):
    add_log(get_label('indexing_test_document'))
    file_name = "test_file.txt"
    content = "This is a test document for ML model analysis."
    return ingest_document(es_client, index, file_name, content)

def process_and_analyze_document(es_client, index, doc_id):
    add_log(get_label('checking_document_processing'))
    try:
        doc = wait_for_document(es_client, index, doc_id)
        if 'ml' not in doc['_source']:
            add_log(get_label('ml_field_not_found'))
            return False, None

        add_log(get_label('ml_field_found'))
        analysis_result, model_info = analyze_ml_field(doc['_source']['ml'])
        add_log(get_label('ml_field_analysis_result'), True)
        add_log(analysis_result, True)

        if not model_info:
            add_log(get_label('no_models_found'))
            return False, None

        return True, model_info
    except TimeoutError as e:
        add_log(get_label('document_processing_timeout', error=str(e)))
        return False, None

def test_watsonx_discovery_connection(config):
    doc_id = None
    try:
        es_client = create_es_client(config)
        index = config['index_name']

        doc_id = create_and_ingest_test_document(es_client, index)
        if not doc_id:
            return False, None

        success, model_info = process_and_analyze_document(es_client, index, doc_id)
        if success:
            st.session_state.model_info = model_info
            st.session_state.wxd_connection_tested = True
            st.session_state.watsonx_discovery_client = es_client
            return True, model_info
        else:
            return False, None

    except Exception as e:
        add_log(get_label('watsonx_discovery_test_error', error=str(e)))
        return False, None
    finally:
        if doc_id:
            add_log(get_label('deleting_test_document'))
            delete_document(es_client, index, doc_id)


def handle_model_selection(model_info):
    st.write(get_label('multiple_models_found'))
    model_options = [get_label('select_option')] + [f"{info['embedding_type']} - {info['model_id']}" for info in model_info]
    selected_model = st.selectbox(get_label('select_model'), model_options, key='model_selection')
    
    if selected_model != get_label('select_option'):
        selected_index = model_options.index(selected_model) - 1
        selected_model_info = model_info[selected_index]
        
        st.write(get_label('selected_model_info', 
                            type=selected_model_info['embedding_type'], 
                            id=selected_model_info['model_id'],
                            field=selected_model_info['embedding_field']))

        st.session_state.watsonx_discovery_config['text_embedding_type'] = selected_model_info
        return True
    else:
        st.write(get_label('no_model_selected'))
        st.session_state.watsonx_discovery_config['text_embedding_type'] = None
        return False

def delete_document(es_client, index, doc_id):
    try:
        result = es_client.delete(index=index, id=doc_id)
        if result['result'] == 'deleted':
            add_log(get_label('document_deletion_success', doc_id=doc_id))
        else:
            add_log(get_label('document_deletion_error', error=str(result)))
    except Exception as e:
        add_log(get_label('document_deletion_exception', error=str(e)), False)

def wait_for_document(es_client, index, doc_id, field_to_check='ml', max_retries=12, delay=10):
    for i in range(max_retries):
        try:
            result = es_client.get(index=index, id=doc_id)
            if field_to_check in result['_source']:
                add_log(get_label('document_field_found', field=field_to_check))
                return result
            add_log(get_label('document_processing_wait', field=field_to_check, attempt=i+1, max_attempts=max_retries), True)
        except Exception as e:
            add_log(get_label('document_wait_error', error=str(e)), False)
        time.sleep(delay)
    raise TimeoutError(get_label('document_processing_timeout'))

# Milvus 설정을 위한 입력 페이지를 생성합니다.
# 사용자로부터 Milvus 연결에 필요한 정보를 입력받습니다.
#
# :return: Milvus 설정 정보를 담은 딕셔너리
def milvus_config_page():
    st.subheader("Milvus 설정")
    
    # Milvus 호스트 및 포트 설정
    milvus_host = st.text_input("Milvus Host", value="")
    milvus_port = st.number_input("Milvus Port", value=8080)
    
    # PEM 파일 업로더 추가
    pem_file = st.file_uploader("Server PEM File", type=['pem'])
    if pem_file:
        # PEM 파일이 업로드된 경우 저장
        server_pem_path = save_pem_file(pem_file)
    else:
        # 업로드된 PEM 파일이 없는 경우 기존 경로 사용
        server_pem_path = st.session_state.get('server_pem_path', '')
    
    # 서버 이름, 사용자 이름, 비밀번호 설정
    server_name = st.text_input("Server Name", value="localhost")
    user = st.text_input("User", value="root")
    password = st.text_input("Password", type="password", value="")
    
    # 컬렉션 이름 설정
    collection_name = st.text_input("Collection Name", value="")

    # Milvus 설정 정보를 딕셔너리로 구성
    milvus_config = {
        "milvus_host": milvus_host,
        "milvus_port": milvus_port,
        "server_pem_path": server_pem_path,
        "server_name": server_name,
        "user": user,
        "password": password,
        "collection_name": collection_name
    }

    return milvus_config

# 업로드된 PEM 파일을 저장하고 파일 경로를 반환합니다.
#
# :param uploaded_file: Streamlit의 파일 업로더를 통해 업로드된 파일 객체
# :return: 저장된 PEM 파일의 경로 또는 None (파일이 없는 경우)
def save_pem_file(uploaded_file):
    if uploaded_file is not None:
        # 임시 디렉토리 생성 (없는 경우)
        temp_dir = os.path.join(os.getcwd(), "temp")
        os.makedirs(temp_dir, exist_ok=True)
        
        # 파일 저장
        file_path = os.path.join(temp_dir, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        # 세션 상태에 파일 경로 저장
        st.session_state.server_pem_path = file_path
        
        return file_path
    return None

# Milvus 연결을 테스트하고 컬렉션의 존재 여부를 확인합니다.
#
#
# :param milvus_config: Milvus 연결 설정 정보를 담은 딕셔너리
# :return: 연결 및 컬렉션 확인 결과 (True 또는 False)
def test_milvus_connection(milvus_config):
    try:
        result = False
        # Milvus 서버에 연결 시도
        if not connect_to_milvus(milvus_config):
            add_log(get_label("milvus_connection_failed_log"), False)
            return result

        if connections.has_connection("default"):
            # 컬렉션 존재 여부 확인
            if utility.has_collection(milvus_config["collection_name"]):
                st.warning(get_label('collection_exists', collection_name=milvus_config['collection_name']))
                result = False
            else:
                add_log(get_label('milvus_connection_success', collection_name=milvus_config['collection_name']), True)
                result = True

            # 연결 해제
            # connections.disconnect("default")
        else:
            add_log(get_label('milvus_connection_failed_log'), False)
            result = False
    except Exception as e:
        # 연결 중 예외 발생 시 오류 로그 추가
        add_log(get_label('connection_error', error=str(e)), False)
        result = False
    return result

def check_and_release_resources():
    if 'last_milvus_use_time' in st.session_state:
        if time.time() - st.session_state.last_milvus_use_time > 3600:  # 1시간 후
            if 'milvus_collection' in st.session_state:
                st.session_state.milvus_collection.release()
            connections.disconnect_all()
            st.session_state.pop('milvus_collection', None)
            st.session_state.pop('last_milvus_use_time', None)



def ensure_milvus_connection():
    try:
        # 컬렉션 로드 상태 확인
        load_state = utility.load_state(st.session_state.milvus_config["collection_name"])
        
        # 로드 상태가 문자열로 반환되므로, 직접 비교
        if load_state != "Loaded":
            st.session_state.milvus_collection.load()
        return True
    except Exception as e:
        add_log(f"Error loading Milvus collection: {str(e)}", False)
        return False



# WatsonX.AI 서비스의 가용성을 테스트합니다.
#
# 간단한 API 요청을 보내 서비스가 정상적으로 응답하는지 확인합니다.
#
# :return: 서비스 가용 여부 (True 또는 False)
def is_available_watsonx():
    if st.session_state.get('login_mode') == 'general':
        return is_available_watsonx_general()
    else:
        return is_available_watsonx_prompt_template()

def check_watsonx_availability(api_url, headers, payload, mode):
    try:
        # API 요청 전송
        response = requests.post(api_url, json=payload, headers=headers)
        
        # 응답 상태 코드 확인
        if 200 <= response.status_code < 300:
            add_log(get_label(f'watsonx_{mode}_connection_success'), True)
            return True
        else:
            add_log(get_label(f'watsonx_{mode}_request_failed', status_code=response.status_code), False)
            add_log(get_label('response_content', content=response.text), False)            
            return False
    except requests.RequestException as e:
        # 요청 중 예외 발생 시 오류 로그 추가
        add_log(get_label(f'watsonx_{mode}_request_error', error=str(e)), False)
        return False

def is_available_watsonx_general():
    api_url = f"{st.session_state.url}{st.session_state.method}"
    headers = {
        "Content-Type": "application/json",
        "Accept": "application/json",
        "Authorization": f"Bearer {st.session_state.bearer}"
    }
    payload = {
        "model_id": "mistralai/mixtral-8x7b-instruct-v01",
        "input": "Hello, World!",
        "parameters": {
            "decoding_method": "greedy",
            "max_new_tokens": 5
        },
        "project_id": st.session_state.project_id
    }
    return check_watsonx_availability(api_url, headers, payload, "general")

def is_available_watsonx_prompt_template():
    api_url = st.session_state.stream_endpoint_url
    headers = {
        "Content-Type": "application/json",
        "Accept": "application/json",
        "Authorization": f"Bearer {st.session_state.bearer}"
    }
    
    prompt_variables = {}
    if st.session_state.get('vector_search_params'):
        prompt_variables[st.session_state.vector_search_params] = "Test vector search value"
    if st.session_state.get('user_input_params'):
        prompt_variables[st.session_state.user_input_params] = "Test user input value"
    
    if not prompt_variables:
        prompt_variables = {"default_question": "Test question"}
    
    payload = {
        "parameters": {
            "prompt_variables": prompt_variables
        }
    }
    return check_watsonx_availability(api_url, headers, payload, "prompt_template")

# 사용자 설정이 완료되었는지 확인합니다.
#
# 기본 설정과 WatsonX.AI 서비스 가용성, 그리고 외부 벡터 검색 방법 사용 여부에 따른 
# 추가 설정을 검사합니다.
#
# :param use_external_vector: 외부 벡터 검색 방법 사용 여부
# :param external_vector_test_success: 외부 벡터 검색 방법 연결 테스트 성공 여부
# :return: 모든 필요한 설정이 완료되었는지 여부 (True 또는 False)
def are_settings_complete(use_external_vector, external_vector_test_success):
    missing_settings = check_basic_settings()
    
    if not missing_settings:
        # WatsonX.AI 서비스 가용성 확인
        watsonx_available = is_available_watsonx()
        
        if not watsonx_available:
            missing_settings.append('WatsonX.AI Availability')
        
        if use_external_vector == VECTOR_SEARCH_FAISS:
            pass  # FAISS 사용 시 추가 확인 불필요
        elif use_external_vector == VECTOR_SEARCH_MILVUS:
            if not external_vector_test_success:
                missing_settings.append('Milvus Connection Test')
        elif use_external_vector == VECTOR_SEARCH_WATSON_DISCOVERY:
            discovery_config = st.session_state.watsonx_discovery_config
            required_keys = ['server_url', 'api_key', 'index_name', 'text_embedding_type']
            
            for key in required_keys:
                value = discovery_config.get(key)
                if key == 'text_embedding_type':
                    if not isinstance(value, dict) or not value:
                        missing_settings.append(f'Watson Discovery {key}')
                else:
                    if not isinstance(value, str) or not value.strip():
                        missing_settings.append(f'Watson Discovery {key}')
            
            if not external_vector_test_success:
                missing_settings.append('Watson Discovery Connection Test')
        else:
            missing_settings.append('Valid Vector Search Method')
    
    return missing_settings


def check_basic_settings():
    missing_settings = []
    
    if st.session_state.get('login_mode') == 'general':
        settings_to_check = [
            ('url', 'URL'),
            ('apikey', 'API Key'),
            ('project_id', 'Project ID'),
            ('bearer', 'Bearer Token'),
            ('selected_model', 'Selected Model'),
        ]
        
        for key, name in settings_to_check:
            if not st.session_state.get(key) or not st.session_state[key].strip():
                missing_settings.append(name)
        
        if not isinstance(st.session_state.get('chunk_size'), int) or st.session_state.chunk_size <= 0:
            missing_settings.append('Valid Chunk Size')
        
        if not isinstance(st.session_state.get('overlap_size'), int) or st.session_state.overlap_size < 0:
            missing_settings.append('Valid Overlap Size')
        
        if st.session_state.get('overlap_size', 0) >= st.session_state.get('chunk_size', 0):
            missing_settings.append('Overlap Size smaller than Chunk Size')
    
    else:  # 프롬프트 템플릿 모드
        settings_to_check = [
            ('stream_endpoint_url', 'Stream Endpoint URL'),
            ('apikey', 'API Key'),
            ('bearer', 'Bearer Token'),
            ('vector_search_params', 'Vector Search Parameters'),
            ('user_input_params', 'User Input Parameters'),
        ]
        
        for key, name in settings_to_check:
            if not st.session_state.get(key) or not st.session_state[key].strip():
                missing_settings.append(name)
        
        if not isinstance(st.session_state.get('chunk_size'), int) or st.session_state.chunk_size <= 0:
            missing_settings.append('Valid Chunk Size')
        
        if not isinstance(st.session_state.get('overlap_size'), int) or st.session_state.overlap_size < 0:
            missing_settings.append('Valid Overlap Size')
        
        if st.session_state.get('overlap_size', 0) >= st.session_state.get('chunk_size', 0):
            missing_settings.append('Overlap Size smaller than Chunk Size')
    
    return missing_settings

    
# FAISS를 사용하여 프롬프트와 가장 유사한 문서 청크를 검색합니다.
#
# :param prompt: 사용자 입력 프롬프트
# :param selected_file: 검색 대상 파일 (모든 파일 또는 특정 파일)
# :return: 검색된 문서 내용과 사용된 파일 목록
def perform_faiss_search(prompt, selected_file, top_k=3):
    query_vector = st.session_state.sentence_transformer.encode([prompt])
    
    file_name = None if selected_file == get_label('all_files') else selected_file
    search_results = search_faiss(query_vector, top_k, file_name)
    
    context = get_label('related_document_content') + "\n"
    used_files = set()
    
    for result in search_results:
        context += get_label('file_content_format', file=result['file_name'], content=result['chunk'])
        used_files.add(result['file_name'])
    
    if not used_files:
        used_files = {get_label('no_search_results_file')}
        context += get_label('no_search_results') + "\n"
    
    return context, used_files

# Milvus를 사용하여 프롬프트와 가장 유사한 문서 청크를 검색합니다.
#
# :param prompt: 사용자 입력 프롬프트
# :param collection: Milvus 컬렉션 객체
# :param selected_file: 검색 대상 파일 (모든 파일 또는 특정 파일)
# :param top_k: 반환할 가장 유사한 문서의 수 (기본값: 3)
# :return: 검색된 문서 내용과 사용된 파일 목록
def perform_milvus_search(prompt, selected_file, top_k=3):

    query_vector = st.session_state.sentence_transformer.encode([prompt])
    context = get_label('related_document_content') + "\n"
    used_files = set()

    try:
        # Collection 객체를 세션 상태에서 가져오거나 생성
        if 'milvus_collection' not in st.session_state:
            st.session_state.milvus_collection = Collection(st.session_state.milvus_config["collection_name"])
            st.session_state.milvus_collection.load()

        collection = st.session_state.milvus_collection

        search_params = {"metric_type": "L2", "params": {"nprobe": 10}}
        
        results = collection.search(
            data=[query_vector[0].tolist()],
            anns_field="vector",
            param=search_params,
            limit=top_k,
            expr=None if selected_file == get_label('all_files') else f'file_name == "{selected_file}"',
            output_fields=["file_name", "content"]
        )

        if not results:
            st.warning(get_label('no_search_results'))
            context += get_label('no_search_results') + "\n"
            used_files.add(get_label('no_search_results_file'))
        else:
            for hits in results:
                for hit in hits:
                    file_name = hit.entity.get('file_name')
                    content = hit.entity.get('content')
                    if file_name and content:
                        context += get_label('file_content_format', file=file_name, content=content)
                        used_files.add(file_name)
                    else:
                        st.warning(get_label('missing_info_in_result', hit=str(hit)))

        if not used_files:
            st.warning(get_label('no_valid_file_info'))

        return context, used_files
    except Exception as e:
        add_log(get_label('milvus_search_error', error=str(e)), False)
        return get_label('search_error_occurred'), set()

# 검색 조건을 만들어 검색을 수행합니다.
#
#:param collection: Milvus 컬렉션 객체
#:param selected_file: 검색 대상 파일 (모든 파일 또는 특정 파일)
#:param top_k: 반환할 가장 유사한 문서의 수 (기본값: 3)
#:param query_vector: 검색을 위한 쿼리 벡터
#:param search_params: 검색 파라미터
#:return: 검색 결과
def make_search_condition(collection, selected_file, top_k, query_vector, search_params):
    if selected_file != get_label('all_files'):
            # 특정 파일에서만 검색
        results = collection.search(
                data=[query_vector[0]],
                anns_field="vector",
                param=search_params,
                limit=top_k,
                expr=f'file_name == "{selected_file}"',
                output_fields=["file_name", "content"]
            )
    else:
            # 모든 파일에서 검색
        results = collection.search(
                data=[query_vector[0]],
                anns_field="vector",
                param=search_params,
                limit=top_k,
                output_fields=["file_name", "content"]
            )
        
    return results

# 채팅 조건을 확인합니다.
#
# :return: 채팅 조건 여부 (True 또는 False)
def is_chat_condition():
    if st.session_state.vector_search_method == VECTOR_SEARCH_FAISS and len(st.session_state.document_vectors) == 0:
        add_log(get_label('upload_file_first'), False)
        return False
    elif st.session_state.vector_search_method == VECTOR_SEARCH_MILVUS and not ensure_milvus_connection():
        st.error(get_label('milvus_connection_failed'))
        return False
    return True  

# 채팅 기록을 표시합니다.
def display_chat_history():
    # 채팅 인터페이스
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.write(message["content"])

def main_page():
    st.markdown("""
    <style>
    .main > div {
        padding-top: 0.5rem;
    }
    .block-container {
        padding-top: 1rem;
        padding-bottom: 0rem;
    }
    </style>
    """, unsafe_allow_html=True)    
    construct_side_bar()
    # Main content
    selected_file = chat_header()
    display_chat_history()
    handle_user_input(selected_file)

# 사용자 입력을 처리합니다.
#
#:param selected_file: 검색 대상 파일 (모든 파일 또는 특정 파일)
def handle_user_input(selected_file):
    if prompt := st.chat_input(get_label('enter_message')):
        display_user_message(prompt)
        # 설정이 완료되었는지 확인
        if is_chat_condition():
            process_chat(selected_file, prompt)

#채팅 프로세스를 처리합니다.
#:param selected_file: 검색 대상 파일 (모든 파일 또는 특정 파일)
#:param prompt: 사용자 입력 프롬프트
def process_chat(selected_file, prompt):
    # 벡터 검색 수행
    context, used_files = vector_search(selected_file, prompt)
    full_prompt = prepare_full_prompt(prompt, context)

    display_ai_response(used_files, full_prompt)

# AI 응답을 표시합니다.
#
#:param used_files: 사용된 파일 목록
#:param full_prompt: 전체 프롬프트
def display_ai_response(used_files, full_prompt):
    # WatsonX.AI API 호출
    with st.chat_message("assistant"):
        full_response = request_message(used_files, full_prompt)

    st.session_state.messages.append({"role": "assistant", "content": full_response})

# API 호출을 위한 조건을 준비합니다.
#
#:param used_files: 사용된 파일 목록
#:param full_prompt: 전체 프롬프트
#:return: API 호출 조건
def request_message(used_files, full_prompt):
    message_placeholder = st.empty()
    full_response = ""
    # API 호출을 위한 매개변수 준비
    parameters = prepare_api_parameters()

    async def stream_response():
        nonlocal full_response
        session = aiohttp.ClientSession()
        try:
            api_url, headers, payload = prepare_request_condition(full_prompt, parameters)
            ensure_valid_token()

            async with session.post(api_url, json=payload, headers=headers) as response:
                if response.status != 200:
                    add_log(get_label('api_call_failed', status=response.status), False)
                    return

                async for line in response.content:
                    if line.startswith(b'data:'):
                        try:
                            data = json.loads(line[5:])
                            generated_text = data['results'][0]['generated_text']
                            full_response += generated_text
                            message_placeholder.markdown(full_response + "▌")
                        except json.JSONDecodeError:
                            pass

        except Exception as e:
            add_log(get_label('error_occurred', error=str(e)), False)
        finally:
            await session.close()

    asyncio.run(stream_response())
    finalize_response(message_placeholder, full_response, used_files)
    return full_response

# 응답 후 처리를 완료합니다.
#
#:param message_placeholder: 메시지 자리표시자
#:param full_response: 전체 응답
#:param used_files: 사용된 파일 목록
def finalize_response(message_placeholder, full_response, used_files):
    message_placeholder.markdown(full_response)
    add_log(get_label('ai_response_log', response=full_response))      
    # 응답 후 확장기 추가
    show_result_expander(used_files)

# API 호출 조건을 준비합니다.
#
#:param full_prompt: 전체 프롬프트
#:param parameters: API 호출 매개변수
#:return: API 호출 조건
def prepare_request_condition(full_prompt, parameters):
    headers = {
        "Content-Type": "application/json",
        "Accept": "application/json",
        "Authorization": f"Bearer {st.session_state.bearer}"
    }

    if st.session_state.get('login_mode') == 'general':
        api_url = f"{st.session_state.url}{st.session_state.method}"
        payload = {
            "model_id": st.session_state.model_id,
            "input": full_prompt,
            "parameters": parameters,
            "project_id": st.session_state.project_id
        }
    else:  # 프롬프트 템플릿 모드
        api_url = st.session_state.stream_endpoint_url
        payload = full_prompt

    return api_url, headers, payload

# 전체 프롬프트를 준비합니다.
#
#:param prompt: 사용자 입력 프롬프트
#:param context: 검색된 문서 내용
#:return: 전체 프롬트
def prepare_full_prompt(prompt, context):
    if st.session_state.get('login_mode') == 'general':
        full_prompt = st.session_state.user_prompt.format(context=context, user_input=prompt)
        st.session_state.last_full_prompt = full_prompt
        add_log(get_label('full_prompt_log', prompt=full_prompt))
        return full_prompt
    else:  # 프롬프트 템플릿 모드
        prompt_variables = {}
        if st.session_state.get('vector_search_params'):
            prompt_variables[st.session_state.vector_search_params] = context
        if st.session_state.get('user_input_params'):
            prompt_variables[st.session_state.user_input_params] = prompt
        
        full_prompt = {
            "parameters": {
                "prompt_variables": prompt_variables
            }
        }
        st.session_state.last_full_prompt = str(full_prompt)  # 로그를 위해 문자열로 변환
        add_log(get_label('full_prompt_log', prompt=str(full_prompt)))
        return full_prompt

#사용자 메시지를 표시합니다.
#:param prompt: 사용자 입력 프롬프트
def display_user_message(prompt):
    st.session_state.messages.append({"role": "user", "content": prompt})
    add_log(get_label('user_log', prompt=prompt))
    with st.chat_message("user"):
        st.write(prompt)

# API 호출 매개변수를 준비합니다.
#:return: API 호출 매개변수
def prepare_api_parameters():
    if st.session_state.get('login_mode') == 'general':
        parameters = {
            GenParams.DECODING_METHOD: st.session_state.decoding_method,
            GenParams.MAX_NEW_TOKENS: st.session_state.max_new_tokens,
            GenParams.MIN_NEW_TOKENS: st.session_state.min_new_tokens,
            GenParams.TEMPERATURE: st.session_state.temperature,
            GenParams.TOP_K: st.session_state.top_k,
            GenParams.TOP_P: st.session_state.top_p
        }
        return parameters
    else:  # 프롬프트 템플릿 모드
        return None

#벡터 검색을 수행합니다.
#:param selected_file: 검색 대상 파일 (모든 파일 또는 특정 파일)
#:param prompt: 사용자 입력 프롬프트
#:return: 검색된 문서 내용과 사용된 파일 목록
def vector_search(selected_file, prompt):
    if st.session_state.vector_search_method == VECTOR_SEARCH_FAISS:
        context, used_files = perform_faiss_search(prompt, selected_file)
    elif st.session_state.vector_search_method == VECTOR_SEARCH_MILVUS:
        context, used_files = perform_milvus_search(prompt, selected_file)
    elif st.session_state.vector_search_method == VECTOR_SEARCH_WATSON_DISCOVERY:
        context, used_files = perform_watsonx_discovery_search(prompt, selected_file)
    else:
        raise ValueError(f"Unknown vector search method: {st.session_state.vector_search_method}")
    return context, used_files

def prepare_search_query(prompt, search_method, model_id, embedding_field, k, num_candidates, min_score):
    if search_method == 'text_embedding':
        return {
            "knn": {
                "field": embedding_field,
                "query_vector_builder": {
                    "text_embedding": {
                        "model_id": model_id,
                        "model_text": prompt
                    }
                },
                "k": k,
                "num_candidates": num_candidates
            },
            "_source": ["file_name", "content"]
        }
    elif search_method == 'text_expanding':
        return {
            "query": {
                "text_expansion": {
                    embedding_field: {
                        "model_id": model_id,
                        "model_text": prompt
                    }
                }
            },
            "size": k,
            "min_score": min_score,
            "_source": ["file_name", "content"]
        }
    else:
        raise ValueError(f"Unknown search method: {search_method}")

def apply_file_filter(query, selected_file):
    if selected_file != get_label('all_files'):
        file_filter = {"term": {"file_name": selected_file}}
        if 'query' in query:
            query['query'] = {"bool": {"must": [query['query'], file_filter]}}
        else:
            query['post_filter'] = file_filter
    return query

def process_search_results(response):
    context = get_label('related_document_content') + "\n"
    used_files = set()
    
    for hit in response['hits']['hits']:
        file_name = hit['_source']['file_name']
        content = hit['_source']['content']
        context += get_label('file_content_format', file=file_name, content=content)
        used_files.add(file_name.split('_chunk_')[0])  # Remove chunk suffix to get original file name
    
    if not used_files:
        used_files = {get_label('no_search_results_file')}
        context += get_label('no_search_results') + "\n"
    
    return context, used_files

def perform_watsonx_discovery_search(prompt, selected_file, k=3, num_candidates=30, min_score=0.5):
    es_client = st.session_state.watsonx_discovery_client
    index = st.session_state.watsonx_discovery_config['index_name']
    embedding_info = st.session_state.watsonx_discovery_config.get('text_embedding_type', {})
    search_method = embedding_info.get('embedding_type', 'text_embedding')
    model_id = embedding_info.get('model_id', '')
    embedding_field = embedding_info.get('embedding_field', '')
    
    try:
        add_log(f"Search Method: {search_method}")
        add_log(f"Model ID: {model_id}")
        add_log(f"Embedding Field: {embedding_field}")
        
        query = prepare_search_query(prompt, search_method, model_id, embedding_field, k, num_candidates, min_score)
        query = apply_file_filter(query, selected_file)
        
        add_log(f"Watson Discovery Query: {json.dumps(query, indent=2)}")
        
        response = es_client.search(index=index, body=query)
        return process_search_results(response)
    
    except Exception as e:
        error_message = get_label('watsonx_discovery_search_error', error=str(e))
        add_log(error_message, False)
        st.error(error_message)
        return get_label('search_error_occurred'), set()

#채팅 헤더를 생성합니다.
#:return: 선택된 파일
def chat_header():
    st.markdown('<div class="header-container">', unsafe_allow_html=True)
    col1, col2 = st.columns([4, 1])
    
    selected_file = None

    with col1:
        st.markdown('<div class="header-left">', unsafe_allow_html=True)
        st.markdown(f'<h1 class="header-text">{get_label("header_title")}</h1>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)  # 추가: header-left div 기
    
    with col2:
        selected_file = st.selectbox(
            get_label("select_file_to_search"),
            [get_label("all_files")] + list(st.session_state.document_vectors.keys()),
            key="file_selector"
        )
        
    st.markdown('<hr style="margin-top: 0.05rem; border-color: rgba(250, 250, 250, 0.2);">', unsafe_allow_html=True)
    if st.toggle(get_label("view_logs")):
        with st.container():
            with st.container():
                st.markdown('<div class="log-container">', unsafe_allow_html=True)
                log_container = st.empty()
                log_container.code(st.session_state.log_messages)
                st.markdown('</div>', unsafe_allow_html=True)    
    if len(st.session_state.document_vectors) < 1:
        st.markdown(f'<div class="toggle-description"><p class="description-text">{get_label("description_text")}</p></div>', unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)  # header-container div 닫기
    
    return selected_file

#사이드바를 구성합니다.
#:return: 사이드바
def construct_side_bar():
    with st.sidebar:
        st.subheader(get_label('document_upload'), help=get_label('document_upload_help'))
        uploaded_files = st.file_uploader(get_label('file_types'), type=['pdf', 'docx', 'xlsx', 'xls', 'pptx'], accept_multiple_files=True)
        
        if uploaded_files:
            process_uploaded_files(uploaded_files)

        st.markdown("---")
        if st.session_state.document_vectors or (st.session_state.vector_search_method == VECTOR_SEARCH_MILVUS and st.session_state.milvus_config):
            st.subheader(get_label('uploaded_files'))
            for filename in st.session_state.document_vectors.keys():
                st.markdown(get_label('file_item', filename=filename))
            st.markdown("---")

        # 일반 모드일 때만 user_prompt와 llm_model_settings를 표시
        if st.session_state.get('login_mode') == 'general':
            st.subheader(get_label('user_prompt'), help=get_label('user_prompt_help'))
            with st.expander(get_label('click_expand_collapse')):
                st.session_state.user_prompt = st.text_area(get_label('prompt_template'), value=st.session_state.user_prompt, height=200)

            st.subheader(get_label('llm_model_settings'), help=get_label('llm_model_settings_help'))
            with st.expander(get_label('click_expand_collapse')):
                llm_model_settings()

        st.sidebar.markdown("---")
        st.sidebar.markdown("### © 2024 IBM Korea Customer Success")
        st.sidebar.markdown("**gwangsu kim**")
        st.sidebar.markdown("[gwangsu.kim@ibm.com](mailto:gwangsu.kim@ibm.com)")


#업로드된 파일을 처리합니다.
#:param uploaded_files: 업로드된 파일 목록
def process_uploaded_files(uploaded_files):
    progress_bar = st.progress(0)
    for i, uploaded_file in enumerate(uploaded_files):
        file_hash = hashlib.md5(uploaded_file.getvalue()).hexdigest()
        
        if file_hash not in st.session_state.processed_files:
            chunks, chunk_vectors = process_uploaded_file_cached(uploaded_file.getvalue(), uploaded_file.name)

            if len(chunks) > 0 and (chunk_vectors is not None or st.session_state.vector_search_method == VECTOR_SEARCH_WATSON_DISCOVERY):
                if st.session_state.vector_search_method == VECTOR_SEARCH_MILVUS:
                    ingest_to_milvus(uploaded_file.name, chunks, chunk_vectors)
                elif st.session_state.vector_search_method == VECTOR_SEARCH_WATSON_DISCOVERY:
                    bulk_ingest_documents(uploaded_file.name, chunks)
                else:  # FAISS
                    add_vectors_to_faiss(uploaded_file.name, chunks, chunk_vectors)
                
                if st.session_state.vector_search_method != VECTOR_SEARCH_WATSON_DISCOVERY:
                    st.session_state.document_vectors[uploaded_file.name] = chunk_vectors
                st.session_state.document_chunks[uploaded_file.name] = chunks
                
                st.session_state.processed_files.add(file_hash)
                add_log(get_label('file_uploaded_vectorized', file_name=uploaded_file.name), True)
            else:
                add_log(get_label('file_processing_failed', file_name=uploaded_file.name), False)
            
            display_vectorization_status()
        
        progress = (i + 1) / len(uploaded_files)
        progress_bar.progress(progress)



def bulk_ingest_documents(file_name, chunks):
    es_client = st.session_state.watsonx_discovery_client
    index = st.session_state.watsonx_discovery_config['index_name']
    try:
        actions = []
        for i, chunk in enumerate(chunks):
            actions.extend([
                {"index": {"_index": index}},
                {
                    "file_name": file_name,
                    "content": chunk,
                    "_extract_binary_content": True,
                    "_reduce_whitespace": True,
                    "_run_ml_inference": True
                }
            ])

        response = es_client.bulk(body=actions, index=index, pipeline=index)

        success = sum(1 for item in response['items'] if item['index']['status'] == 201)
        failed = len(response['items']) - success

        add_log(get_label('bulk_indexing_success', success_count=success, file_name=file_name))
        
        if failed:
            add_log(get_label('bulk_indexing_partial_failure', failed_count=failed, file_name=file_name))
            for item in response['items']:
                if item['index']['status'] != 201:
                    add_log(get_label('bulk_indexing_failure_detail', error=str(item['index']['error'])))

        return success, failed

    except Exception as e:
        add_log(get_label('bulk_indexing_exception', error=str(e), file_name=file_name))
        return 0, len(chunks)

def initialize_faiss_index(vector_dimension):
    if 'faiss_index' not in st.session_state:
        st.session_state.faiss_index = faiss.IndexFlatL2(vector_dimension)
        st.session_state.faiss_id_to_file_chunk = {}
        add_log(get_label('faiss_index_initialized', dimension=vector_dimension), True)

def add_vectors_to_faiss(file_name, chunks, vectors):
    if 'faiss_index' not in st.session_state or st.session_state.faiss_index is None:
        # FAISS 인덱스가 없으면 새로 생성
        vector_dimension = vectors.shape[1]
        st.session_state.faiss_index = faiss.IndexFlatL2(vector_dimension)
        add_log(get_label('faiss_index_initialized', dimension=vector_dimension), True)
    
    start_id = st.session_state.faiss_index.ntotal
    st.session_state.faiss_index.add(vectors)
    
    for i, chunk in enumerate(chunks):
        st.session_state.faiss_id_to_file_chunk[start_id + i] = (file_name, chunk)
    
    add_log(get_label('vectors_added_to_faiss', count=len(chunks), file_name=file_name), True)

def search_faiss(query_vector, top_k=3, file_name=None):
    if 'faiss_index' not in st.session_state or st.session_state.faiss_index.ntotal == 0:
        add_log(get_label('faiss_index_empty'), False)
        return []

    if file_name:
        # 특정 파일에 해당하는 인덱스만 추출
        file_indices = [idx for idx, (fname, _) in st.session_state.faiss_id_to_file_chunk.items() if fname == file_name]
        if not file_indices:
            add_log(get_label('file_not_found_in_index', file_name=file_name), False)
            return []
        
        # 추출된 인덱스에 대해서만 검색 수행
        subset_index = faiss.IndexFlatL2(st.session_state.faiss_index.d)
        vectors_to_search = st.session_state.faiss_index.reconstruct_n(file_indices[0], len(file_indices))
        subset_index.add(vectors_to_search)
        distances, subset_indices = subset_index.search(query_vector.reshape(1, -1), min(top_k, len(file_indices)))
        indices = [file_indices[i] for i in subset_indices[0] if i != -1]
    else:
        # 전체 인덱스에 대해 검색 수행
        distances, indices = st.session_state.faiss_index.search(query_vector.reshape(1, -1), top_k)
        indices = indices[0]

    results = []
    for distance, idx in zip(distances[0], indices):
        if idx != -1:  # FAISS returns -1 for not found
            file_name, chunk = st.session_state.faiss_id_to_file_chunk[idx]
            results.append({
                "distance": float(distance),
                "file_name": file_name,
                "chunk": chunk
            })
    
    return results



#벡터 검색 설정을 표시합니다.
def vector_search_settings():
    st.info(get_label('embedding_model_info'))
    st.text(get_label('vector_search_method', method=st.session_state.vector_search_method))
    
    if st.session_state.vector_search_method != VECTOR_SEARCH_WATSON_DISCOVERY:
        col1, col2, col3 = st.columns([2, 1, 1])
        
        with col1:
            st.session_state.selected_model = st.selectbox(
                get_label('embedding_model'),
                st.session_state.config["model"]["embedding_models"],
                index=0
            )
    else:
        col2, col3 = st.columns(2)
    
    with col2:
        st.session_state.chunk_size = st.number_input(
            get_label('chunk_size'), 
            value=st.session_state.chunk_size, 
            min_value=1
        )
    
    with col3:
        st.session_state.overlap_size = st.number_input(
            get_label('overlap_size'), 
            value=st.session_state.overlap_size, 
            min_value=0
        )

#LLM 모델 설정을 표시합니다.
def llm_model_settings():
    st.session_state.model_id = st.selectbox(get_label('model_id'), st.session_state.config["model"]["supported_models"], index=0)
    st.session_state.decoding_method = st.selectbox(get_label('decoding_method'), ["greedy", "sample"], index=0)
    st.session_state.max_new_tokens = st.number_input(get_label('max_new_tokens'), value=st.session_state.max_new_tokens, min_value=1)
    st.session_state.min_new_tokens = st.number_input(get_label('min_new_tokens'), value=st.session_state.min_new_tokens, min_value=1)
    st.session_state.temperature = st.slider(get_label('temperature'), min_value=0.0, max_value=1.0, value=st.session_state.temperature, step=0.1)
    st.session_state.top_k = st.number_input(get_label('top_k'), value=st.session_state.top_k, min_value=1)
    st.session_state.top_p = st.number_input(get_label('top_p'), value=st.session_state.top_p, min_value=0.0, max_value=1.0, step=0.1)

#사용된 파일과 프롬프트를 표시합니다.
#:param used_files: 사용된 파일 목록
def show_result_expander(used_files):
    # 확장기 표시 로직
    expander = st.expander(get_label('view_files_and_prompt'), expanded=False)
    with expander:
        st.subheader(get_label('files_used_in_search'))
        for filename in used_files:
            st.markdown(get_label('file_item', filename=filename))
        
        st.subheader(get_label('prompt'))
        st.text_area(get_label('full_prompt'), value=st.session_state.last_full_prompt, height=300, disabled=True)
    
    # 확장기를 세션 상태에 저장
    st.session_state.expanders.append(expander)


def log_basic_faiss_info(faiss_index):
    total_vectors = faiss_index.ntotal
    dimension = faiss_index.d

    add_log(get_label("total_vectors_in_index", count=total_vectors))
    add_log(get_label("vector_dimension", dimension=dimension))

    if hasattr(faiss_index, 'nlist'):
        add_log(get_label("number_of_clusters", count=faiss_index.nlist))

    index_type = type(faiss_index).__name__
    add_log(get_label("faiss_index_type", type=index_type))

    if isinstance(faiss_index, faiss.IndexFlatL2):
        add_log(get_label("index_flat_l2_info"))
    elif isinstance(faiss_index, faiss.IndexIVFFlat):
        add_log(get_label("index_ivf_flat_info", nlist=faiss_index.nlist, nprobe=faiss_index.nprobe))
    elif isinstance(faiss_index, faiss.IndexIVFPQ):
        add_log(get_label("index_ivf_pq_info", nlist=faiss_index.nlist, nprobe=faiss_index.nprobe, M=faiss_index.pq.M, nbits=faiss_index.pq.nbits))

    memory_usage = total_vectors * dimension * 4 / (1024 * 1024)  # MB 단위
    add_log(get_label("estimated_memory_usage", usage=f"{memory_usage:.2f} MB"))

def log_file_chunk_info():
    if 'faiss_id_to_file_chunk' in st.session_state:
        file_chunk_map = st.session_state.faiss_id_to_file_chunk
        unique_files = len(set(file_name for file_name, _ in file_chunk_map.values()))
        add_log(get_label("unique_files_in_index", count=unique_files))

        file_chunk_counts = {}
        for file_name, _ in file_chunk_map.values():
            file_chunk_counts[file_name] = file_chunk_counts.get(file_name, 0) + 1

        add_log(get_label("chunks_per_file"))
        for file_name, count in file_chunk_counts.items():
            add_log(get_label("file_chunks", file=file_name, count=count))

def log_faiss_status():
    add_log(get_label("faiss_status"))

    try:
        if 'faiss_index' not in st.session_state or st.session_state.faiss_index is None:
            add_log(get_label("faiss_index_not_initialized"))
            return

        faiss_index = st.session_state.faiss_index
        log_basic_faiss_info(faiss_index)
        log_file_chunk_info()

    except AttributeError as e:
        add_log(get_label("faiss_data_not_found", error=str(e)))
    except Exception as e:
        add_log(get_label("faiss_logging_error", error=str(e)))


def get_collection_details(collection_name):
    try:
        collection = Collection(collection_name)
        
        add_log(get_label("collection_details", name=collection_name))
        
        schema = collection.schema
        add_log(get_label("schema_info"))
        add_log(get_label("description", desc=schema.description))
        add_log(get_label("fields"))
        for field in schema.fields:
            add_log(get_label("field_info", name=field.name, type=field.dtype, is_primary=field.is_primary))
        
        add_log(get_label("basic_info"))
        add_log(get_label("name", name=collection.name))
        add_log(get_label("description", desc=collection.description))
        add_log(get_label("is_empty", empty=collection.is_empty))
        add_log(get_label("entity_count", count=collection.num_entities))
        
        primary_field = collection.primary_field
        add_log(get_label("primary_key_info"))
        if primary_field:
            add_log(get_label("field_info", name=primary_field.name, type=primary_field.dtype, is_primary=True))
        else:
            add_log(get_label("no_primary_key"))
        
        partitions = collection.partitions
        add_log(get_label("partition_info"))
        for partition in partitions:
            add_log(get_label("partition_details", name=partition.name, desc=partition.description))
        
        indexes = collection.indexes
        add_log(get_label("index_info"))
        if indexes:
            for index in indexes:
                add_log(get_label("field_name", name=index.field_name))
                add_log(get_label("index_name", name=index.index_name))
                add_log(get_label("index_type", type=index.params.get('index_type', 'N/A')))
                add_log(get_label("metric_type", type=index.params.get('metric_type', 'N/A')))
                add_log(get_label("additional_params", params=index.params.get('params', {})))
        else:
            add_log(get_label("no_index"))
        
        load_state = utility.load_state(collection_name)
        add_log(get_label("load_state", state=load_state))
        
    except Exception as e:
        add_log(get_label("error_occurred", error=str(e)))

#벡터화 상태를 표시합니다.
def display_vectorization_status():
    add_log("##### " + get_label("vectorization_status_title", method=st.session_state.vector_search_method) + " #####")
    
    if st.session_state.vector_search_method == VECTOR_SEARCH_FAISS:
        log_faiss_status()    
    elif st.session_state.vector_search_method == VECTOR_SEARCH_MILVUS:
        get_collection_details(st.session_state.milvus_config["collection_name"])
    elif st.session_state.vector_search_method == VECTOR_SEARCH_WATSON_DISCOVERY:
        log_watsonx_discovery_status()
    
    add_log(get_label("vectorized_files"))
    for file in st.session_state.document_vectors.keys():
        add_log(f"- {file}")
        
    add_log(get_label("status_separator"))

def log_watsonx_discovery_status():
    try:
        es_client = st.session_state.watsonx_discovery_client
        index_name = st.session_state.watsonx_discovery_config['index_name']
        
        # 인덱스 정보 가져오기
        index_info = es_client.indices.get(index=index_name)
        
        # 문서 수 가져오기
        doc_count = es_client.count(index=index_name)['count']
        
        add_log(get_label("watsonx_discovery_index_name", name=index_name))
        add_log(get_label("watsonx_discovery_doc_count", count=doc_count))
        
        # 인덱스 매핑 정보 (필드 정보) 표시
        if 'mappings' in index_info[index_name]:
            add_log(get_label("watsonx_discovery_fields"))
            for field, field_info in index_info[index_name]['mappings']['properties'].items():
                add_log(f"  - {field}: {field_info.get('type', 'Unknown type')}")
        
    except Exception as e:
        add_log(get_label("watsonx_discovery_status_error", error=str(e)))


###################################################
def main():
    # 초기화 여부를 확인하는 플래그를 session_state에 추가
    if 'authenticated' not in st.session_state:
        st.session_state.authenticated = False

    # 초기화가 아직 수행되지 않았다면 initialize() 실행
    if not st.session_state.authenticated:
        initialize()
        init_session_status()
        login_page()
    else:       
        main_page()
###################################################


if __name__ == "__main__":
    main()
